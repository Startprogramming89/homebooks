import streamlit as st
import os
from pathlib import Path
import streamlit as st
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, Settings
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI
import requests
from ebooklib import epub
from bs4 import BeautifulSoup
import PyPDF2
import docx
from pptx import Presentation

# Page configuration
st.set_page_config(
    page_title="Document Chatbot",
    page_icon="📚",
    layout="wide"
)

# Initialize session state for the index
if 'index' not in st.session_state:
    st.session_state.index = None
    st.session_state.processing_complete = False

def init_openai():
    """Initialize OpenAI settings"""
    api_key = st.secrets["OPENAI_API_KEY"]
    if not api_key:
        st.error("❌ OpenAI API key not found in .env file!")
        st.stop()
    return api_key

def get_files_from_repo(owner, repo, branch='main'):
    """Get all supported files from a GitHub repository"""
    api_url = f"https://api.github.com/repos/{owner}/{repo}/git/trees/{branch}?recursive=1"
    response = requests.get(api_url)
    response.raise_for_status()
    
    supported_extensions = {'.epub', '.pdf', '.docx', '.doc', '.pptx', '.ppt'}
    files = []
    found_files = {ext: [] for ext in supported_extensions}
    
    for item in response.json().get('tree', []):
        ext = Path(item['path']).suffix.lower()
        if ext in supported_extensions:
            download_url = f"https://raw.githubusercontent.com/{owner}/{repo}/{branch}/{item['path']}"
            files.append((item['path'], download_url))
            found_files[ext].append(item['path'])
    
    return files, found_files

def download_from_github(download_url, local_filename):
    """Download a file from GitHub"""
    response = requests.get(download_url)
    response.raise_for_status()
    
    with open(local_filename, 'wb') as f:
        f.write(response.content)
    return local_filename

def read_epub(file_path):
    """Process EPUB files"""
    book = epub.read_epub(file_path)
    text = ""
    for item in book.get_items():
        if isinstance(item, epub.EpubHtml):
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            text += soup.get_text() + "\n"
    return text

def read_pdf(file_path):
    """Process PDF files"""
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    return text

def read_docx(file_path):
    """Process DOCX files"""
    doc = docx.Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def read_pptx(file_path):
    """Process PPTX files"""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def process_file(file_path):
    """Process a file based on its extension"""
    ext = Path(file_path).suffix.lower()
    try:
        if ext == '.epub':
            return read_epub(file_path)
        elif ext == '.pdf':
            return read_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return read_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return read_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    except Exception as e:
        st.error(f"Error processing {file_path}: {str(e)}")
        return None

def process_documents(files):
    """Process documents and create index"""
    try:
        # Create documents directory if it doesn't exist
        if not os.path.exists("documents"):
            os.makedirs("documents")
        
        processed_files = []
        
        for filename, download_url in files:
            try:
                local_path = os.path.join("documents", os.path.basename(filename))
                st.info(f"📥 Downloading: {filename}")
                
                # Download the file
                download_from_github(download_url, local_path)
                
                # Process file and save text
                st.info(f"📖 Processing: {filename}")
                text = process_file(local_path)
                
                if text:
                    # Save processed text
                    txt_filename = f"{os.path.splitext(local_path)[0]}.txt"
                    with open(txt_filename, "w", encoding="utf-8") as f:
                        f.write(text)
                    
                    processed_files.append(txt_filename)
                    st.success(f"✅ Successfully processed: {filename}")
                
            except Exception as e:
                st.error(f"❌ Error processing {filename}: {str(e)}")
        
        if processed_files:
            documents = SimpleDirectoryReader(input_files=processed_files).load_data()
            Settings.embed_model = OpenAIEmbedding()
            Settings.llm = OpenAI(model="gpt-4o-mini-2024-07-18")
            index = VectorStoreIndex.from_documents(documents)
            st.session_state.index = index
            st.session_state.processing_complete = True
            return True
        return False
    except Exception as e:
        st.error(f"Error processing documents: {str(e)}")
        return False

# Main UI
st.title("📚 Document Chatbot")

# Sidebar
with st.sidebar:
    st.header("Configuration")
    github_owner = st.text_input("GitHub Username", value="Startprogramming89")
    github_repo = st.text_input("GitHub Repository", value="homebooks")
    github_branch = st.text_input("Branch", value="main")
    
    if st.button("Process Documents"):
        with st.spinner("Processing documents..."):
            # Initialize OpenAI
            init_openai()
            
            # Get and display files from repository
            try:
                files, found_files = get_files_from_repo(github_owner, github_repo, github_branch)
                
                # Display found files
                st.success("Files found in repository:")
                for ext, file_list in found_files.items():
                    if file_list:
                        st.write(f"\n{ext.upper()} files ({len(file_list)}):")
                        for file in file_list:
                            st.write(f"  - {file}")
                
                # Process documents with the files list
                if process_documents(files):
                    st.success("✅ Documents processed successfully!")
                else:
                    st.error("❌ No documents were processed.")
                    
            except Exception as e:
                st.error(f"Error: {str(e)}")

# Main chat interface
if st.session_state.processing_complete:
    st.header("Chat with your documents")
    
    # Initialize chat history
    if "messages" not in st.session_state:
        st.session_state.messages = []

    # Display chat history
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Example prompts section
    st.markdown("### 📝 Example Questions You Can Ask:")
    example_prompts = [
        "🇩🇪 German Examples:",
        "• Bitte finde und gib den Originalabschnitt oder das Kapitel aus dem Buch \"Babyjahre Entwicklung und Erziehung in den ersten vier Jahren\" zurück, das das Verhalten eines Kindes erklärt, das Dinge auf den Boden wirft. Bitte fasse nichts zusammen, ich möchte den vollständigen Originaltext.",
        
        "🇬🇧 English Examples:",
        "• Please find and return the original paragraph or chapter in the book \"xxx\" that explains \"the behavior of a child throwing things\". Please do not summarize - I want the full original wording.",
        "• What are the conclusions drawn in these documents?",
        "• Find any references to specific dates or timelines.",
        "• Explain the methodology described in these documents."
    ]

    # Display examples in an expander
    with st.expander("Click to see example questions"):
        for prompt in example_prompts:
            st.markdown(prompt)

    # Chat input
    if prompt := st.chat_input("Ask a question about your documents"):
        # Add user message to chat history
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # Generate response
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                query_engine = st.session_state.index.as_query_engine()
                response = query_engine.query(prompt)
                st.markdown(response.response)
                st.session_state.messages.append({"role": "assistant", "content": response.response})

else:
    st.info("👈 Please configure your GitHub repository and process documents using the sidebar.")
